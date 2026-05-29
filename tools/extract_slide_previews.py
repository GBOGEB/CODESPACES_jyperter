#!/usr/bin/env python3
"""
Extract slide text previews and styled digital-twin hints from a PPTX and write:
- slide_texts.json (full texts per slide)
- slide_digital_twins.json (compact summary + style hints)
- slide_preview.html (simple visual preview with inline CSS)

Usage:
  python tools/extract_slide_previews.py \
    --pptx output/handover_final/QPLANT_comparison.pptx \
    --outdir output/handover_final \
    --limit 20
"""
import argparse
import html
import json
import re
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as exc:
    raise SystemExit("python-pptx is required. pip install python-pptx") from exc


def extract_slides(pptx_path: Path, limit: int | None = None):
    """Return extracted text, counts, and style hints for each slide."""
    prs = Presentation(str(pptx_path))
    slides = []
    all_slides = list(prs.slides)
    rng = all_slides[:limit] if limit else all_slides
    for i, s in enumerate(rng):
        texts = []
        bullets = 0
        for sh in s.shapes:
            if hasattr(sh, 'text'):
                t = (sh.text or '').strip()
                if t:
                    texts.append(t)
                    if any(t.lstrip().startswith(ch) for ch in ('-', '•', '*')):
                        bullets += 1
        title = texts[0] if texts else ''
        # NOTE: `isupper()` returns False for strings with no cased characters
        # (e.g. all-digit or all-punctuation titles), so some legitimate title
        # slides may be styled as body slides. This is an intentional rough
        # heuristic; use placeholder_format.type from python-pptx for stricter
        # detection if needed.
        is_title = bool(title) and len(title) < 120 and title.isupper()
        bg = '#f0f8ff' if is_title else ('#ffffff' if bullets > 0 else '#fafafa')
        slides.append({
            'n': i + 1,
            'texts': texts,
            'shape_count': len(s.shapes),
            'bullets': bullets,
            'style_hint': {
                'bg': bg,
                'is_title': is_title,
            },
        })
    return slides


_HEX_COLOR_RE = re.compile(r'^#[0-9a-fA-F]{3}(?:[0-9a-fA-F]{3})?$')
_FALLBACK_BG = '#ffffff'


def _safe_bg(color: str) -> str:
    """Return *color* if it matches a hex-color literal, else a safe fallback."""
    return color if _HEX_COLOR_RE.match(color) else _FALLBACK_BG


def write_html_preview(slides: list[dict], out_html: Path):
    """Write a small HTML preview for extracted slide content."""
    css = """
    <style>
      body{font-family:Arial,Helvetica,sans-serif;background:#111;color:#eee;margin:0;padding:24px}
      .deck{display:flex;flex-wrap:wrap;gap:16px}
      .card{width:360px;height:240px;border-radius:10px;box-shadow:0 2px 12px rgba(0,0,0,.4);padding:12px;overflow:auto}
      .title{margin:4px 0 6px 0;color:#111}
      .meta{font-size:12px;color:#444;margin-bottom:6px}
      .para{margin:4px 0;color:#222;font-size:12px}
      .num{background:#444;color:#fff;border-radius:6px;padding:2px 6px;font-size:11px}
    </style>
    """
    parts = ["<html><head><meta charset='utf-8'><title>Slide Preview</title>", css, "</head><body>"]
    parts.append("<h2>Slide Preview (first {} slides)</h2>".format(len(slides)))
    parts.append("<div class='deck'>")
    for s in slides:
        bg = _safe_bg(s.get('style_hint', {}).get('bg', _FALLBACK_BG))
        texts = s.get('texts', [])
        title = texts[0] if texts else ''
        parts.append("<div class='card' style='background:{}'>".format(bg))
        parts.append("<div class='meta'><span class='num'>#{} </span> • shapes:{} • bullets:{}</div>".format(s['n'], s['shape_count'], s['bullets']))
        if title:
            parts.append("<h3 class='title'>{}</h3>".format(html.escape(title[:180])))
        for t in texts[1:5]:
            parts.append("<div class='para'>{}</div>".format(html.escape(t[:220])))
        parts.append("</div>")
    parts.append("</div></body></html>")
    out_html.write_text("\n".join(parts), encoding='utf-8')


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--pptx', required=True)
    ap.add_argument('--outdir', required=True)
    ap.add_argument('--limit', type=int, default=20)
    args = ap.parse_args()

    pptx = Path(args.pptx)
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    slides = extract_slides(pptx, args.limit)

    (outdir / 'slide_texts.json').write_text(json.dumps(slides, indent=2, ensure_ascii=False), encoding='utf-8')
    twins = [
        {
            'slide': s['n'],
            'summary': s['texts'][:3],
            'style_hint': s['style_hint'],
            'shape_count': s['shape_count'],
            'bullets': s['bullets'],
        }
        for s in slides
    ]
    (outdir / 'slide_digital_twins.json').write_text(json.dumps(twins, indent=2, ensure_ascii=False), encoding='utf-8')
    write_html_preview(slides, outdir / 'slide_preview.html')
    print('WROTE:', outdir / 'slide_texts.json')
    print('WROTE:', outdir / 'slide_digital_twins.json')
    print('WROTE:', outdir / 'slide_preview.html')


if __name__ == '__main__':
    main()
