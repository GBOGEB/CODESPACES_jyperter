
"""
PowerPoint Parser (.pptx)
=========================

Parser for Microsoft PowerPoint files using python-pptx.
Extracts slides, text, images, and presentation structure.
"""

from typing import Dict, Any, List
import logging

from .base_parser import BaseParser

logger = logging.getLogger(__name__)

class PowerPointParser(BaseParser):
    """Parser for PowerPoint .pptx files"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx']
        
    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        Parse PowerPoint file and extract content
        
        Args:
            file_path: Path to .pptx file
            
        Returns:
            Dictionary with presentation analysis
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            # Extract slides
            slides_data = []
            total_text_shapes = 0
            total_images = 0
            all_text = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_info = self._extract_slide_data(slide, slide_idx)
                slides_data.append(slide_info)
                
                total_text_shapes += slide_info["text_shapes_count"]
                total_images += slide_info["images_count"]
                all_text.extend(slide_info["text_content"])
            
            # Extract presentation properties
            properties = self._extract_properties(prs)
            
            # Analyze presentation structure
            structure = self._analyze_structure(slides_data)
            
            return {
                "slide_count": len(prs.slides),
                "total_text_shapes": total_text_shapes,
                "total_images": total_images,
                "word_count": len(" ".join(all_text).split()),
                "slides": slides_data,
                "properties": properties,
                "structure": structure,
                "presentation_type": self._identify_presentation_type(all_text, slides_data)
            }
            
        except ImportError:
            logger.error("python-pptx library not installed. Install with: pip install python-pptx")
            return {"error": "python-pptx library not available"}
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file {file_path}: {e}")
            return {"error": str(e)}
    
    def extract_text(self, file_path: str) -> str:
        """
        Extract plain text from PowerPoint file
        
        Args:
            file_path: Path to .pptx file
            
        Returns:
            Plain text content from all slides
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            for slide_idx, slide in enumerate(prs.slides):
                text_content.append(f"=== Slide {slide_idx + 1} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            
            return "\n".join(text_content)
            
        except ImportError:
            return "Error: python-pptx library not available"
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint file {file_path}: {e}")
            return f"Error: {e}"
    
    def _extract_slide_data(self, slide, slide_idx: int) -> Dict[str, Any]:
        """Extract data from a single slide"""
        slide_data = {
            "slide_number": slide_idx + 1,
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
            "text_content": [],
            "text_shapes_count": 0,
            "images_count": 0,
            "tables_count": 0,
            "charts_count": 0,
            "shapes_info": []
        }
        
        for shape in slide.shapes:
            shape_info = {
                "shape_type": str(shape.shape_type),
                "has_text": hasattr(shape, "text")
            }
            
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_data["text_content"].append(shape.text)
                slide_data["text_shapes_count"] += 1
                shape_info["text"] = shape.text
            
            # Count different shape types
            if shape.shape_type.name == 'PICTURE':
                slide_data["images_count"] += 1
            elif shape.shape_type.name == 'TABLE':
                slide_data["tables_count"] += 1
                shape_info["table_data"] = self._extract_table_from_shape(shape)
            elif shape.shape_type.name == 'CHART':
                slide_data["charts_count"] += 1
            
            slide_data["shapes_info"].append(shape_info)
        
        return slide_data
    
    def _extract_table_from_shape(self, shape) -> Dict[str, Any]:
        """Extract table data from a table shape"""
        try:
            if hasattr(shape, 'table'):
                table = shape.table
                rows = []
                
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text.strip())
                    rows.append(cells)
                
                return {
                    "row_count": len(table.rows),
                    "column_count": len(table.columns),
                    "data": rows
                }
        except Exception as e:
            logger.debug(f"Could not extract table data: {e}")
        
        return {"error": "Could not extract table data"}
    
    def _extract_properties(self, prs) -> Dict[str, Any]:
        """Extract presentation properties"""
        properties = {}
        
        try:
            core_props = prs.core_properties
            properties.update({
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "keywords": core_props.keywords or ""
            })
        except Exception as e:
            logger.debug(f"Could not extract presentation properties: {e}")
        
        return properties
    
    def _analyze_structure(self, slides_data: List[Dict]) -> Dict[str, Any]:
        """Analyze presentation structure"""
        structure = {
            "layouts_used": [],
            "text_distribution": [],
            "content_balance": {
                "text_heavy_slides": 0,
                "image_heavy_slides": 0,
                "balanced_slides": 0
            }
        }
        
        layouts = set()
        
        for slide in slides_data:
            layouts.add(slide["layout_name"])
            
            # Analyze content balance
            text_count = slide["text_shapes_count"]
            image_count = slide["images_count"]
            
            if text_count > image_count * 2:
                structure["content_balance"]["text_heavy_slides"] += 1
            elif image_count > text_count * 2:
                structure["content_balance"]["image_heavy_slides"] += 1
            else:
                structure["content_balance"]["balanced_slides"] += 1
            
            structure["text_distribution"].append({
                "slide": slide["slide_number"],
                "text_shapes": text_count,
                "word_count": len(" ".join(slide["text_content"]).split())
            })
        
        structure["layouts_used"] = list(layouts)
        return structure
    
    def _identify_presentation_type(self, all_text: List[str], 
                                  slides_data: List[Dict]) -> str:
        """Identify presentation type based on content"""
        text_content = " ".join(all_text).lower()
        
        # Check for common presentation types
        if any(word in text_content for word in ['agenda', 'meeting', 'discussion']):
            return "Meeting Presentation"
        elif any(word in text_content for word in ['training', 'tutorial', 'learn']):
            return "Training/Educational"
        elif any(word in text_content for word in ['sales', 'product', 'demo']):
            return "Sales/Demo"
        elif any(word in text_content for word in ['report', 'results', 'analysis']):
            return "Report/Analysis"
        elif any(word in text_content for word in ['proposal', 'plan', 'strategy']):
            return "Proposal/Strategy"
        else:
            return "General Presentation"
