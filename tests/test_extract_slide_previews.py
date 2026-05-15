import importlib.util
from pathlib import Path

from pptx import Presentation


MODULE_PATH = (
    Path(__file__).resolve().parents[1] / "tools" / "extract_slide_previews.py"
)
SPEC = importlib.util.spec_from_file_location("extract_slide_previews", MODULE_PATH)
MODULE = importlib.util.module_from_spec(SPEC)
assert SPEC and SPEC.loader
SPEC.loader.exec_module(MODULE)

extract_slides = MODULE.extract_slides
write_html_preview = MODULE.write_html_preview


def test_extract_and_preview(tmp_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "<script>TITLE</script>"
    slide.placeholders[1].text = "- bullet1\n<script>bullet2</script>"
    pptx_path = tmp_path / "sample.pptx"
    prs.save(pptx_path)

    slides = extract_slides(pptx_path)
    assert slides[0]["texts"][0] == "<script>TITLE</script>"
    html_path = tmp_path / "preview.html"
    write_html_preview(slides, html_path)
    assert html_path.exists()
    content = html_path.read_text()
    assert "Slide Preview" in content
    assert "<script>" not in content
    assert "&lt;script&gt;TITLE&lt;/script&gt;" in content
